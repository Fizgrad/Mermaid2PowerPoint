#!/usr/bin/env python3
"""
Render a Mermaid diagram into a single-slide PowerPoint deck.

Dependencies (install separately):
- Mermaid CLI: npm i -g @mermaid-js/mermaid-cli
- python-pptx: pip install -r requirements.txt
"""
from __future__ import annotations

import argparse
import shutil
import subprocess
import sys
import tempfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Iterable, Tuple


EMU_PER_INCH = 914400
EMU_PER_PX_AT_96_DPI = EMU_PER_INCH / 96


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Convert a Mermaid diagram to a one-slide PPTX with the diagram centered.",
        formatter_class=argparse.ArgumentDefaultsHelpFormatter,
    )
    parser.add_argument(
        "mermaid_file",
        help="Path to the .mmd file. Use '-' to read the Mermaid definition from stdin.",
    )
    parser.add_argument(
        "-o",
        "--output",
        default="diagram.pptx",
        help="Output PPTX path.",
    )
    parser.add_argument(
        "--slide-width",
        type=float,
        default=10.0,
        help="Slide width in inches.",
    )
    parser.add_argument(
        "--slide-height",
        type=float,
        default=7.5,
        help="Slide height in inches.",
    )
    parser.add_argument(
        "--margin",
        type=float,
        default=0.5,
        help="Margin (in inches) to leave around the rendered diagram.",
    )
    parser.add_argument(
        "--theme",
        default="default",
        help="Mermaid theme passed to mermaid-cli.",
    )
    parser.add_argument(
        "--background",
        default="white",
        help="Background color passed to mermaid-cli (CSS color or hex).",
    )
    parser.add_argument(
        "--scale",
        type=float,
        default=1.0,
        help="Scale factor passed to mermaid-cli when rendering.",
    )
    parser.add_argument(
        "--format",
        choices=["png", "emf"],
        default="png",
        help="Image format embedded into the slide. Use 'emf' for vector graphics you can ungroup/edit in PPT.",
    )
    parser.add_argument(
        "--mmdc-path",
        help="Path to the mermaid-cli binary (mmdc). If omitted, the script looks it up on PATH.",
    )
    parser.add_argument(
        "--inkscape-path",
        help="Path to the Inkscape binary (needed only for EMF export).",
    )
    parser.add_argument(
        "--keep-intermediate",
        action="store_true",
        help="Do not delete the intermediate rendered PNG (useful for debugging).",
    )
    return parser.parse_args()


def ensure_dependencies() -> None:
    try:
        import pptx  # noqa: F401
    except ImportError as exc:
        sys.exit(
            "python-pptx is required. Install dependencies with: pip install -r requirements.txt"
        )
    if not find_mermaid_cli():
        sys.exit(
            "mermaid-cli (mmdc) not found. Install it with: npm i -g @mermaid-js/mermaid-cli"
        )


def find_mermaid_cli() -> str | None:
    candidates: Iterable[str] = (
        ["mmdc", "node_modules/.bin/mmdc"]
    )
    for candidate in candidates:
        resolved = shutil.which(candidate)
        if resolved:
            return resolved
    return None


def read_mermaid_input(mermaid_file: str) -> Path:
    if mermaid_file != "-":
        return Path(mermaid_file).expanduser().resolve()

    tmp_file = Path(tempfile.mkstemp(prefix="mermaid-", suffix=".mmd")[1])
    tmp_file.write_text(sys.stdin.read(), encoding="utf-8")
    return tmp_file


def render_mermaid(
    mermaid_path: Path,
    output_path: Path,
    *,
    mmdc_path: str,
    theme: str,
    background: str,
    scale: float,
    output_format: str = "png",
) -> None:
    cmd = [
        mmdc_path,
        "-i",
        str(mermaid_path),
        "-o",
        str(output_path),
        "-t",
        theme,
        "-b",
        background,
        "-e",
        output_format,
    ]
    if scale and scale != 1.0:
        cmd.extend(["-s", str(scale)])

    try:
        subprocess.run(cmd, check=True)
    except FileNotFoundError:
        sys.exit(
            "Could not run mermaid-cli (mmdc). Ensure it is installed and on PATH."
        )
    except subprocess.CalledProcessError as exc:
        sys.exit(f"Mermaid rendering failed: {exc}")


def convert_svg_to_emf(svg_path: Path, emf_path: Path, *, inkscape_path: str) -> None:
    cmd = [
        inkscape_path,
        str(svg_path),
        "--export-type=emf",
        f"--export-filename={emf_path}",
    ]
    try:
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    except FileNotFoundError:
        sys.exit(
            "Inkscape not found. Install it or specify --inkscape-path when using --format emf."
        )
    except subprocess.CalledProcessError as exc:
        sys.exit(f"Inkscape EMF export failed: {exc.stderr.decode().strip()}")


def _parse_length_px(raw: str | None) -> float | None:
    if not raw:
        return None
    raw = raw.strip()
    for suffix in ("px", "pt"):
        if raw.endswith(suffix):
            raw = raw[: -len(suffix)]
            break
    try:
        return float(raw)
    except ValueError:
        return None


def svg_size_px(svg_path: Path) -> Tuple[int, int]:
    tree = ET.parse(svg_path)
    root = tree.getroot()
    width = _parse_length_px(root.attrib.get("width"))
    height = _parse_length_px(root.attrib.get("height"))
    if width is None or height is None:
        viewbox = root.attrib.get("viewBox", "")
        parts = viewbox.split()
        if len(parts) == 4:
            try:
                width = float(parts[2])
                height = float(parts[3])
            except ValueError:
                pass
    if width is None or height is None:
        sys.exit("Could not determine SVG dimensions for scaling.")
    return int(width), int(height)


def image_size_px(image_path: Path) -> Tuple[int, int]:
    try:
        from PIL import Image
    except ImportError:
        sys.exit("Pillow is required (pulled in by python-pptx). Please reinstall dependencies.")

    with Image.open(image_path) as img:
        return img.size


def image_size_emu(image_path: Path) -> Tuple[int, int]:
    width_px, height_px = image_size_px(image_path)
    width_emu = int(width_px * EMU_PER_PX_AT_96_DPI)
    height_emu = int(height_px * EMU_PER_PX_AT_96_DPI)
    return width_emu, height_emu


def add_image_to_slide(
    pptx_path: Path,
    image_path: Path,
    *,
    slide_width_in: float,
    slide_height_in: float,
    margin_in: float,
    image_px: Tuple[int, int] | None = None,
) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(slide_width_in)
    presentation.slide_height = Inches(slide_height_in)
    blank_layout = presentation.slide_layouts[6]  # blank slide
    slide = presentation.slides.add_slide(blank_layout)

    if image_px:
        img_width_emu = int(image_px[0] * EMU_PER_PX_AT_96_DPI)
        img_height_emu = int(image_px[1] * EMU_PER_PX_AT_96_DPI)
    else:
        img_width_emu, img_height_emu = image_size_emu(image_path)
    margin = Inches(max(margin_in, 0))
    available_width = presentation.slide_width - margin * 2
    available_height = presentation.slide_height - margin * 2

    scale = min(
        available_width / img_width_emu,
        available_height / img_height_emu,
        1.0,
    )
    scaled_width = int(img_width_emu * scale)
    scaled_height = int(img_height_emu * scale)

    left = (presentation.slide_width - scaled_width) // 2
    top = (presentation.slide_height - scaled_height) // 2

    # python-pptx expects a filesystem path string, not a Path-like object.
    slide.shapes.add_picture(
        str(image_path), left, top, width=scaled_width, height=scaled_height
    )
    presentation.save(pptx_path)


def main() -> None:
    args = parse_args()
    ensure_dependencies()

    mmdc_path = args.mmdc_path or find_mermaid_cli()
    if not mmdc_path:
        sys.exit(
            "mermaid-cli (mmdc) not found. Install it with: npm i -g @mermaid-js/mermaid-cli"
        )
    if args.format == "emf":
        inkscape = args.inkscape_path or shutil.which("inkscape")
        if not inkscape:
            sys.exit(
                "Inkscape is required for --format emf. Install it or pass --inkscape-path."
            )
    else:
        inkscape = None

    mermaid_source = read_mermaid_input(args.mermaid_file)
    pptx_path = Path(args.output).expanduser().resolve()

    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir_path = Path(tmpdir)
        if args.format == "emf":
            tmp_svg = tmpdir_path / "diagram.svg"
            tmp_emf = tmpdir_path / "diagram.emf"
            render_mermaid(
                mermaid_source,
                tmp_svg,
                mmdc_path=mmdc_path,
                theme=args.theme,
                background=args.background,
                scale=args.scale,
                output_format="svg",
            )
            convert_svg_to_emf(tmp_svg, tmp_emf, inkscape_path=inkscape)
            svg_px = svg_size_px(tmp_svg)
            add_image_to_slide(
                pptx_path,
                tmp_emf,
                slide_width_in=args.slide_width,
                slide_height_in=args.slide_height,
                margin_in=args.margin,
                image_px=svg_px,
            )
            if args.keep_intermediate:
                kept_svg = pptx_path.with_suffix(".svg")
                kept_emf = pptx_path.with_suffix(".emf")
                shutil.copy(tmp_svg, kept_svg)
                shutil.copy(tmp_emf, kept_emf)
                print(f"Kept intermediate vector files at: {kept_svg}, {kept_emf}")
        else:
            tmp_png = tmpdir_path / "diagram.png"
            render_mermaid(
                mermaid_source,
                tmp_png,
                mmdc_path=mmdc_path,
                theme=args.theme,
                background=args.background,
                scale=args.scale,
            )
            add_image_to_slide(
                pptx_path,
                tmp_png,
                slide_width_in=args.slide_width,
                slide_height_in=args.slide_height,
                margin_in=args.margin,
            )

            if args.keep_intermediate:
                kept_png = pptx_path.with_suffix(".png")
                shutil.copy(tmp_png, kept_png)
                print(f"Kept intermediate render at: {kept_png}")

    print(f"✅ Wrote {pptx_path}")


if __name__ == "__main__":
    main()
